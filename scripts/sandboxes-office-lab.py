#!/usr/bin/env python3
"""Live validation for an Office custom disk image in ACA Sandboxes."""

from __future__ import annotations

import hashlib
import json
import os
import subprocess
import sys
from pathlib import Path

try:
    from azure.containerapps.sandbox import (
        AutoDeletePolicy,
        AutoSuspendPolicy,
        EgressPolicy,
        LifecyclePolicy,
        RegistryCredentials,
        SandboxGroupClient,
        endpoint_for_region,
    )
    from azure.identity import AzureCliCredential
except ImportError as error:
    raise SystemExit(
        "Install the preview SDK first: pip3 install azure-containerapps-sandbox"
    ) from error


def run_az(*args: str) -> str:
    result = subprocess.run(
        ["az", *args],
        check=True,
        capture_output=True,
        text=True,
        timeout=600,
    )
    return result.stdout.strip()


def acr_credentials(acr_name: str) -> RegistryCredentials:
    token = json.loads(
        run_az("acr", "login", "--name", acr_name, "--expose-token", "-o", "json")
    )
    return RegistryCredentials(
        username=token["username"],
        token=token["refreshToken"],
    )


def main() -> None:
    subscription_id = os.environ.get("SUBSCRIPTION_ID") or run_az(
        "account", "show", "--query", "id", "--output", "tsv"
    )
    resource_group = os.environ.get(
        "RESOURCE_GROUP", "rg-ai-workspace-sandbox-lab"
    )
    location = os.environ.get("LOCATION", "koreacentral")
    sandbox_group = os.environ.get(
        "SANDBOX_GROUP_NAME", "ai-workspace-sandboxes"
    )
    acr_name = os.environ.get(
        "ACR_NAME", f"aiwssbx{subscription_id.replace('-', '')[:20]}"
    )
    repository = os.environ.get("IMAGE_REPOSITORY", "office-sandbox")
    work_dir = Path(
        os.environ.get("SANDBOX_OFFICE_WORK_DIR", ".work/sandboxes-live")
    )
    work_dir.mkdir(parents=True, exist_ok=True)

    credential = AzureCliCredential()
    client = SandboxGroupClient(
        endpoint_for_region(location),
        credential,
        subscription_id=subscription_id,
        resource_group=resource_group,
        sandbox_group=sandbox_group,
    )
    sandbox = None
    evidence: dict[str, object] = {
        "disk_image_auth": "existing disk image",
    }
    try:
        image_tag = os.environ.get("IMAGE_TAG")
        disk_images = [
            item
            for item in client.list_disk_images()
            if item.labels.get("name", "").startswith("office-")
            and item.status
            and item.status.state in {"Ready", "Succeeded"}
        ]
        if image_tag:
            disk_image_name = f"office-{image_tag}"
            disk_image = next(
                (
                    item
                    for item in disk_images
                    if item.labels.get("name") == disk_image_name
                ),
                None,
            )
        else:
            disk_image = max(
                disk_images,
                key=lambda item: item.labels["name"],
                default=None,
            )
        if disk_image is None:
            image_tag = image_tag or run_az(
                "acr",
                "repository",
                "show-tags",
                "--name",
                acr_name,
                "--repository",
                repository,
                "--orderby",
                "time_desc",
                "--top",
                "1",
                "--query",
                "[0]",
                "--output",
                "tsv",
            )
            if not image_tag:
                raise RuntimeError(
                    f"No ready Office disk image or ACR tag found for "
                    f"{acr_name}/{repository}"
                )
            image = f"{acr_name}.azurecr.io/{repository}:{image_tag}"
            disk_image_name = f"office-{image_tag}"
            print(f"Building disk image from {image}...", flush=True)
            disk_image = client.begin_create_disk_image(
                image,
                name=disk_image_name,
                entrypoint=["/bin/sh", "-c"],
                cmd=["sleep infinity"],
                registry_credentials=acr_credentials(acr_name),
                polling_timeout=1200,
                polling_interval=10,
            ).result()
            evidence["source_image"] = image
            evidence["disk_image_auth"] = "temporary ACR refresh token"
        else:
            evidence["source_image"] = str(disk_image.image.base)
        evidence["disk_image_id"] = disk_image.id
        evidence["disk_image_state"] = disk_image.status.state

        print("Creating Office sandbox...", flush=True)
        sandbox = client.begin_create_sandbox(
            disk=None,
            disk_id=disk_image.id,
            cpu="2000m",
            memory="4096Mi",
            auto_suspend_seconds=1800,
            auto_suspend_mode="Memory",
            egress_policy=EgressPolicy(
                default_action="Deny",
                traffic_inspection="Full",
            ),
            labels={"lab": "office-live-validation"},
        ).result()
        sandbox.set_lifecycle_policy(
            LifecyclePolicy(
                auto_suspend=AutoSuspendPolicy(
                    enabled=True,
                    interval=300,
                    mode="Memory",
                ),
                auto_delete=AutoDeletePolicy(
                    enabled=True,
                    delete_interval_seconds=3600,
                ),
            )
        )
        evidence["sandbox_id"] = sandbox.sandbox_id

        health = sandbox.exec(
            "python --version && libreoffice --version && "
            "pandoc --version | head -1 && pdftotext -v 2>&1 | head -1 && "
            "python -c 'import docx,pptx,openpyxl; "
            'print("office-python-libs=ok")\''
        )
        if (
            health.exit_code != 0
            or "LibreOffice" not in health.stdout
            or "office-python-libs=ok" not in health.stdout
        ):
            raise RuntimeError(f"Office health check failed: {health.stderr}")
        evidence["health"] = health.stdout.strip()

        generation = r"""import hashlib, json, pathlib, subprocess
from docx import Document
from pptx import Presentation
import openpyxl

out=pathlib.Path('/work')
out.mkdir(parents=True, exist_ok=True)
title='AI Workspace ACA Sandboxes validation'
content='Office custom disk image live validation - review draft'

doc=Document()
doc.add_heading(title, level=1)
doc.add_paragraph(content)
doc.save(out/'report.docx')

prs=Presentation()
slide=prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text=title
slide.placeholders[1].text=content
prs.save(out/'slides.pptx')

wb=openpyxl.Workbook()
ws=wb.active
ws.title='Summary'
ws['A1']=title
ws['A2']=content
wb.save(out/'report.xlsx')

for source in ['report.docx','slides.pptx']:
    result=subprocess.run(
        ['libreoffice','--headless','--convert-to','pdf','--outdir',str(out),str(out/source)],
        capture_output=True,
        text=True,
    )
    if result.returncode:
        raise RuntimeError(f'{source} conversion failed: {result.stdout} {result.stderr}')

files={}
for name in ['report.docx','report.pdf','slides.pptx','slides.pdf','report.xlsx']:
    data=(out/name).read_bytes()
    files[name]={'size':len(data),'sha256':hashlib.sha256(data).hexdigest()}
print(json.dumps(files,sort_keys=True))
"""
        sandbox.write_file("/work/generate.py", generation)
        generated = sandbox.exec("python /work/generate.py")
        if generated.exit_code != 0:
            raise RuntimeError(
                f"Office generation failed: {generated.stdout}\n{generated.stderr}"
            )
        files = json.loads(generated.stdout)
        for name, metadata in files.items():
            payload = sandbox.read_file(f"/work/{name}")
            if (
                len(payload) != metadata["size"]
                or hashlib.sha256(payload).hexdigest()
                != metadata["sha256"]
            ):
                raise RuntimeError(f"{name} hash validation failed")
            if name.endswith(".pdf") and not payload.startswith(b"%PDF"):
                raise RuntimeError(f"{name} is not a PDF")
        evidence["generated_files"] = files

        edit = r"""import hashlib, json
from docx import Document
from pptx import Presentation
import openpyxl

old='review draft'
new='review approved'
doc=Document('/work/report.docx')
for paragraph in doc.paragraphs:
    for run in paragraph.runs:
        run.text=run.text.replace(old,new)
doc.save('/work/report.docx')

prs=Presentation('/work/slides.pptx')
for slide in prs.slides:
    for shape in slide.shapes:
        if getattr(shape,'has_text_frame',False):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text=run.text.replace(old,new)
prs.save('/work/slides.pptx')

wb=openpyxl.load_workbook('/work/report.xlsx')
ws=wb.active
ws.title='Final'
ws['B2']='approved-draft'
wb.save('/work/report.xlsx')

result={}
for name in ['report.docx','slides.pptx','report.xlsx']:
    data=open('/work/'+name,'rb').read()
    result[name]=hashlib.sha256(data).hexdigest()
print(json.dumps(result,sort_keys=True))
"""
        sandbox.write_file("/work/edit.py", edit)
        edited = sandbox.exec("python /work/edit.py")
        if edited.exit_code != 0:
            raise RuntimeError(f"Office editing failed: {edited.stderr}")
        edit_hashes = json.loads(edited.stdout)
        for name, digest in edit_hashes.items():
            if (
                hashlib.sha256(sandbox.read_file(f"/work/{name}")).hexdigest()
                != digest
            ):
                raise RuntimeError(f"{name} edit hash validation failed")
        evidence["edited_files"] = edit_hashes

        pdf_text = sandbox.exec("pdftotext /work/slides.pdf - | head -20")
        if pdf_text.exit_code != 0 or "AI Workspace" not in pdf_text.stdout:
            raise RuntimeError("PDF text extraction failed")
        evidence["pdf_text_verified"] = True

        egress = sandbox.exec(
            "python -c \"import urllib.request;"
            "\ntry:"
            "\n urllib.request.urlopen('https://example.com',timeout=5);"
            " print('UNEXPECTED_EGRESS_ALLOWED')"
            "\nexcept Exception as e:"
            "\n print('EGRESS_BLOCKED',type(e).__name__)\""
        )
        if (
            "EGRESS_BLOCKED" not in egress.stdout
            or "UNEXPECTED_EGRESS_ALLOWED" in egress.stdout
        ):
            raise RuntimeError(f"Egress was not blocked: {egress.stdout}")
        evidence["egress"] = egress.stdout.strip()

        before = hashlib.sha256(
            sandbox.read_file("/work/report.docx")
        ).hexdigest()
        sandbox.begin_stop().result()
        stopped = sandbox.get().state
        sandbox.begin_resume().result()
        resumed = sandbox.get().state
        after = hashlib.sha256(
            sandbox.read_file("/work/report.docx")
        ).hexdigest()
        if stopped != "Stopped" or resumed != "Running" or before != after:
            raise RuntimeError("Office suspend/resume did not preserve state")
        evidence["suspend_resume"] = {
            "stopped": stopped,
            "resumed": resumed,
            "state_preserved": True,
        }
        evidence["result"] = "passed"
        print("Office Sandbox live validation passed.", flush=True)
    finally:
        if sandbox is not None:
            try:
                sandbox.delete()
            except Exception as error:
                print(
                    f"WARNING: Sandbox cleanup will rely on auto-delete: {error}",
                    file=sys.stderr,
                )
            finally:
                sandbox.close()
        client.close()
        (work_dir / "office-validation.json").write_text(
            json.dumps(evidence, indent=2, ensure_ascii=False) + "\n",
            encoding="utf-8",
        )


if __name__ == "__main__":
    try:
        main()
    except (RuntimeError, subprocess.CalledProcessError) as error:
        print(f"ERROR: {error}", file=sys.stderr)
        raise SystemExit(1) from error
