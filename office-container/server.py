import hashlib
import json
import mimetypes
import re
import shutil
import subprocess
import tempfile
import threading
import time
import uuid
from http import HTTPStatus
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from urllib.parse import unquote, urlparse

import openpyxl
import pptx
import docx
from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.utils.exceptions import IllegalCharacterError
from pptx import Presentation


BASE_DIR = Path("/work")
MAX_REQUEST_BYTES = 1024 * 1024
MAX_TITLE_CHARS = 200
MAX_CONTENT_CHARS = 100_000
MAX_JOBS = 20
MAX_STORAGE_BYTES = 256 * 1024 * 1024
JOB_TTL_SECONDS = 3600
MAX_EDIT_OPERATIONS = 50
EXCEL_CELL_MAX_CHARS = 32_767
PORT = 8080
JOB_LOCK = threading.Lock()
JOB_OPERATION_LOCK = threading.RLock()

GENERATED_FILES = frozenset(
    {"report.docx", "report.pdf", "report.pptx", "report.xlsx"}
)

# 허용 변환 조합만 처리한다. 임의 source-target 조합은 거부한다.
CONVERSION_MATRIX: dict[tuple[str, str], str] = {
    ("report.docx", "pdf"): "report.docx.pdf",
    ("report.pptx", "pdf"): "report.pptx.pdf",
    ("report.xlsx", "pdf"): "report.xlsx.pdf",
    ("report.docx", "txt"): "report.docx.txt",
}

# 선언적 편집 operation만 허용한다. shell이나 LibreOffice argument는 받지 않는다.
ALLOWED_EDIT_OPERATIONS = frozenset({"setCell", "renameSheet", "replaceText"})

CELL_REFERENCE = re.compile(r"^[A-Z]{1,3}[1-9][0-9]{0,6}$")
SHEET_NAME = re.compile(r"^[A-Za-z0-9 _가-힣-]{1,31}$")
ILLEGAL_XML_CONTROL = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F]")


def tool_version(command: list[str]) -> str:
    result = subprocess.run(
        command,
        check=True,
        capture_output=True,
        text=True,
        timeout=10,
    )
    output = result.stdout.strip() or result.stderr.strip()
    return output.splitlines()[0]


TOOL_VERSIONS = {
    "libreoffice": tool_version(["libreoffice", "--version"]),
    "openpyxl": openpyxl.__version__,
    "pandoc": tool_version(["pandoc", "--version"]),
    "pdftotext": tool_version(["pdftotext", "-v"]),
    "python-docx": docx.__version__,
    "python-pptx": pptx.__version__,
}

EDITED_FILES = frozenset(
    {"report.edited.docx", "report.edited.xlsx"}
)
ALLOWED_DOWNLOADS = (
    GENERATED_FILES | frozenset(CONVERSION_MATRIX.values()) | EDITED_FILES
)


def file_metadata(path: Path, job_id: str) -> dict[str, object]:
    digest = hashlib.sha256()
    with path.open("rb") as source:
        for chunk in iter(lambda: source.read(64 * 1024), b""):
            digest.update(chunk)
    return {
        "name": path.name,
        "size": path.stat().st_size,
        "sha256": digest.hexdigest(),
        "downloadPath": f"/files/{job_id}/{path.name}",
    }


def directory_size(path: Path) -> int:
    return sum(item.stat().st_size for item in path.rglob("*") if item.is_file())


def cleanup_expired_jobs(now: float | None = None) -> None:
    cutoff = (now or time.time()) - JOB_TTL_SECONDS
    for job_dir in BASE_DIR.iterdir():
        if job_dir.is_dir() and job_dir.stat().st_mtime < cutoff:
            shutil.rmtree(job_dir)


class EditError(ValueError):
    """선언적 편집 operation이 유효하지 않은 경우."""


def _edit_set_cell(job_dir: Path, operation: dict[str, object]) -> Path:
    sheet_name = operation.get("sheet")
    cell = operation.get("cell")
    value = operation.get("value")
    if not isinstance(cell, str) or not CELL_REFERENCE.match(cell):
        raise EditError(f"Invalid cell reference: {cell!r}")
    if not isinstance(value, (str, int, float)):
        raise EditError("value must be a string or number")
    if isinstance(value, str):
        if len(value) > 1000:
            raise EditError("value must be at most 1000 characters")
        if ILLEGAL_XML_CONTROL.search(value):
            raise EditError("value contains unsupported control characters")
        # 수식 주입을 막는다. 필요한 수식은 별도 허용 목록으로 다룬다.
        if value.startswith(("=", "+", "-", "@")):
            raise EditError("Formula values are not allowed")

    source = job_dir / "report.edited.xlsx"
    if not source.is_file():
        source = job_dir / "report.xlsx"
    if not source.is_file():
        raise EditError("report.xlsx not found in job")

    workbook = load_workbook(source)
    if sheet_name is None:
        worksheet = workbook.active
    else:
        if not isinstance(sheet_name, str) or sheet_name not in workbook.sheetnames:
            raise EditError(f"Sheet not found: {sheet_name!r}")
        worksheet = workbook[sheet_name]
    worksheet[cell] = value
    output = job_dir / "report.edited.xlsx"
    workbook.save(output)
    return output


def _edit_rename_sheet(job_dir: Path, operation: dict[str, object]) -> Path:
    source_name = operation.get("sheet")
    new_name = operation.get("name")
    if not isinstance(new_name, str) or not SHEET_NAME.match(new_name):
        raise EditError(f"Invalid sheet name: {new_name!r}")

    source = job_dir / "report.edited.xlsx"
    if not source.is_file():
        source = job_dir / "report.xlsx"
    if not source.is_file():
        raise EditError("report.xlsx not found in job")

    workbook = load_workbook(source)
    if source_name is None:
        worksheet = workbook.active
    else:
        if not isinstance(source_name, str) or source_name not in workbook.sheetnames:
            raise EditError(f"Sheet not found: {source_name!r}")
        worksheet = workbook[source_name]
    worksheet.title = new_name
    output = job_dir / "report.edited.xlsx"
    workbook.save(output)
    return output


def _edit_replace_text(job_dir: Path, operation: dict[str, object]) -> Path:
    find = operation.get("find")
    replace = operation.get("replace")
    if not isinstance(find, str) or not find:
        raise EditError("find must be a non-empty string")
    if not isinstance(replace, str):
        raise EditError("replace must be a string")
    if len(find) > 200 or len(replace) > 1000:
        raise EditError("find or replace text is too long")
    if ILLEGAL_XML_CONTROL.search(find) or ILLEGAL_XML_CONTROL.search(replace):
        raise EditError("find or replace contains unsupported control characters")

    source = job_dir / "report.edited.docx"
    if not source.is_file():
        source = job_dir / "report.docx"
    if not source.is_file():
        raise EditError("report.docx not found in job")

    document = Document(str(source))
    for paragraph in document.paragraphs:
        for run in paragraph.runs:
            if find in run.text:
                run.text = run.text.replace(find, replace)
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        if find in run.text:
                            run.text = run.text.replace(find, replace)
    output = job_dir / "report.edited.docx"
    document.save(str(output))
    return output


EDIT_HANDLERS = {
    "setCell": _edit_set_cell,
    "renameSheet": _edit_rename_sheet,
    "replaceText": _edit_replace_text,
}


def apply_edit_operations(
    job_dir: Path, operations: list[dict[str, object]]
) -> list[Path]:
    """Apply a complete edit batch on temporary copies, then commit it."""
    with tempfile.TemporaryDirectory(prefix="edit-", dir=job_dir) as temporary:
        temporary_dir = Path(temporary)
        for name in (
            "report.docx",
            "report.xlsx",
            "report.edited.docx",
            "report.edited.xlsx",
        ):
            source = job_dir / name
            if source.is_file():
                shutil.copyfile(source, temporary_dir / name)

        outputs: dict[str, Path] = {}
        for operation in operations:
            handler = EDIT_HANDLERS[str(operation["op"])]
            result = handler(temporary_dir, operation)
            outputs[result.name] = result

        backups: dict[str, Path | None] = {}
        for name in outputs:
            destination = job_dir / name
            if destination.is_file():
                backup = temporary_dir / f"backup-{name}"
                shutil.copyfile(destination, backup)
                backups[name] = backup
            else:
                backups[name] = None

        committed = []
        try:
            for name in sorted(outputs):
                destination = job_dir / name
                outputs[name].replace(destination)
                committed.append(name)
        except OSError as commit_error:
            rollback_errors = []
            for name in reversed(committed):
                destination = job_dir / name
                backup = backups[name]
                try:
                    if backup is None:
                        destination.unlink(missing_ok=True)
                    else:
                        backup.replace(destination)
                except OSError as rollback_error:
                    rollback_errors.append(rollback_error)
            if rollback_errors:
                raise OSError("Edit commit and rollback both failed") from commit_error
            raise
        return [job_dir / name for name in committed]


def spreadsheet_literal(value: str) -> str:
    """Store formula-like user input as text in generated workbooks."""
    return f"'{value}" if value.startswith(("=", "+", "-", "@")) else value


def remove_job(job_dir: Path) -> None:
    try:
        shutil.rmtree(job_dir)
    except OSError as error:
        print(
            json.dumps(
                {
                    "level": "error",
                    "message": "Failed to remove job directory",
                    "jobDirectory": str(job_dir),
                    "error": str(error),
                }
            ),
            flush=True,
        )


class Handler(BaseHTTPRequestHandler):
    server_version = "OfficeSandbox/1.0"

    def send_json(self, status: HTTPStatus, payload: dict[str, object]) -> None:
        data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)

    def do_GET(self) -> None:
        parsed = urlparse(self.path)
        if parsed.path == "/health":
            self.send_json(
                HTTPStatus.OK,
                {
                    "status": "ok",
                    "tools": TOOL_VERSIONS,
                    "operations": {
                        "generate": ["docx", "pdf", "pptx", "xlsx"],
                        "convert": [
                            {"source": source, "target": target}
                            for source, target in sorted(CONVERSION_MATRIX)
                        ],
                        "edit": sorted(ALLOWED_EDIT_OPERATIONS),
                    },
                    "limits": {
                        "maxRequestBytes": MAX_REQUEST_BYTES,
                        "maxJobs": MAX_JOBS,
                        "maxStorageBytes": MAX_STORAGE_BYTES,
                        "maxEditOperations": MAX_EDIT_OPERATIONS,
                        "jobTtlSeconds": JOB_TTL_SECONDS,
                    },
                },
            )
            return

        parts = [unquote(part) for part in parsed.path.split("/") if part]
        if len(parts) == 3 and parts[0] == "files":
            with JOB_OPERATION_LOCK:
                job_id, filename = parts[1], parts[2]
                if not job_id.isalnum() or filename not in ALLOWED_DOWNLOADS:
                    self.send_json(
                        HTTPStatus.BAD_REQUEST,
                        {"error": "Invalid file path"},
                    )
                    return

                path = BASE_DIR / job_id / filename
                if not path.is_file():
                    self.send_json(HTTPStatus.NOT_FOUND, {"error": "File not found"})
                    return

                self.send_response(HTTPStatus.OK)
                self.send_header(
                    "Content-Type",
                    mimetypes.guess_type(filename)[0] or "application/octet-stream",
                )
                self.send_header("Content-Length", str(path.stat().st_size))
                self.send_header(
                    "Content-Disposition",
                    f'attachment; filename="{filename}"',
                )
                self.end_headers()
                with path.open("rb") as source:
                    shutil.copyfileobj(source, self.wfile, length=64 * 1024)
            return

        self.send_json(HTTPStatus.NOT_FOUND, {"error": "Not found"})

    def _read_json_body(self) -> dict[str, object] | None:
        """공통 JSON 본문 파서. 오류 응답까지 처리하고 None을 반환한다."""
        content_type = self.headers.get("Content-Type", "")
        if not content_type.lower().startswith("application/json"):
            self.send_json(
                HTTPStatus.UNSUPPORTED_MEDIA_TYPE,
                {"error": "Content-Type must be application/json"},
            )
            return None

        try:
            content_length = int(self.headers.get("Content-Length", "0"))
        except ValueError:
            self.send_json(HTTPStatus.BAD_REQUEST, {"error": "Invalid Content-Length"})
            return None

        if content_length <= 0 or content_length > MAX_REQUEST_BYTES:
            self.send_json(
                HTTPStatus.REQUEST_ENTITY_TOO_LARGE,
                {"error": "Request body must be between 1 byte and 1 MB"},
            )
            return None

        try:
            payload = json.loads(self.rfile.read(content_length))
        except (json.JSONDecodeError, UnicodeDecodeError):
            self.send_json(HTTPStatus.BAD_REQUEST, {"error": "Invalid JSON"})
            return None
        if not isinstance(payload, dict):
            self.send_json(
                HTTPStatus.BAD_REQUEST,
                {"error": "JSON body must be an object"},
            )
            return None
        return payload

    def _resolve_job(self, payload: dict[str, object]) -> Path | None:
        """요청의 jobId를 검증하고 job 디렉터리를 돌려준다."""
        job_id = payload.get("jobId")
        if not isinstance(job_id, str) or not job_id.isalnum() or len(job_id) > 64:
            self.send_json(HTTPStatus.BAD_REQUEST, {"error": "Invalid jobId"})
            return None
        job_dir = BASE_DIR / job_id
        if not job_dir.is_dir():
            self.send_json(HTTPStatus.NOT_FOUND, {"error": "Job not found"})
            return None
        return job_dir

    def do_POST(self) -> None:
        path = urlparse(self.path).path
        with JOB_OPERATION_LOCK:
            if path == "/generate":
                self.handle_generate()
                return
            if path == "/convert":
                self.handle_convert()
                return
            if path == "/edit":
                self.handle_edit()
                return
        self.send_json(HTTPStatus.NOT_FOUND, {"error": "Not found"})

    def handle_convert(self) -> None:
        """허용 matrix에 있는 source-target 조합만 변환한다."""
        payload = self._read_json_body()
        if payload is None:
            return
        job_dir = self._resolve_job(payload)
        if job_dir is None:
            return

        source = payload.get("source")
        target = payload.get("target")
        if not isinstance(source, str) or not isinstance(target, str):
            self.send_json(
                HTTPStatus.BAD_REQUEST,
                {"error": "source and target are required strings"},
            )
            return

        output_name = CONVERSION_MATRIX.get((source, target))
        if output_name is None:
            self.send_json(
                HTTPStatus.BAD_REQUEST,
                {
                    "error": "Conversion is not allowed",
                    "allowed": [
                        {"source": allowed_source, "target": allowed_target}
                        for allowed_source, allowed_target in sorted(CONVERSION_MATRIX)
                    ],
                },
            )
            return

        source_path = job_dir / source
        if not source_path.is_file():
            self.send_json(HTTPStatus.NOT_FOUND, {"error": f"{source} not found"})
            return

        job_id = job_dir.name
        output_path = job_dir / output_name
        with tempfile.TemporaryDirectory(
            prefix=f"convert-{job_id}-",
            dir=job_dir,
        ) as temporary:
            temporary_dir = Path(temporary)
            produced = temporary_dir / f"{source_path.stem}.{target}"
            try:
                subprocess.run(
                    [
                        "libreoffice",
                        f"-env:UserInstallation=file:///tmp/lo-profile-{job_id}",
                        "--headless",
                        "--convert-to",
                        target,
                        "--outdir",
                        str(temporary_dir),
                        str(source_path),
                    ],
                    check=True,
                    capture_output=True,
                    text=True,
                    timeout=120,
                )
            except subprocess.TimeoutExpired:
                self.send_json(
                    HTTPStatus.GATEWAY_TIMEOUT,
                    {"error": "Conversion timed out"},
                )
                return
            except subprocess.CalledProcessError as error:
                self.send_json(
                    HTTPStatus.INTERNAL_SERVER_ERROR,
                    {
                        "error": "Conversion failed",
                        "stderr": (error.stderr or "")[-2000:],
                    },
                )
                return

            if not produced.is_file():
                self.send_json(
                    HTTPStatus.INTERNAL_SERVER_ERROR,
                    {"error": "Converted artifact was not produced"},
                )
                return
            shutil.copyfile(produced, output_path)

        self.send_json(
            HTTPStatus.OK,
            {
                "jobId": job_id,
                "source": source,
                "target": target,
                "files": [file_metadata(output_path, job_id)],
            },
        )

    def handle_edit(self) -> None:
        """선언적 편집 operation만 적용한다. shell argument는 받지 않는다."""
        payload = self._read_json_body()
        if payload is None:
            return
        job_dir = self._resolve_job(payload)
        if job_dir is None:
            return

        operations = payload.get("operations")
        if not isinstance(operations, list) or not operations:
            self.send_json(
                HTTPStatus.BAD_REQUEST,
                {"error": "operations must be a non-empty array"},
            )
            return
        if len(operations) > MAX_EDIT_OPERATIONS:
            self.send_json(
                HTTPStatus.BAD_REQUEST,
                {"error": f"operations must be at most {MAX_EDIT_OPERATIONS}"},
            )
            return
        for operation in operations:
            if not isinstance(operation, dict):
                self.send_json(
                    HTTPStatus.BAD_REQUEST,
                    {"error": "each operation must be an object"},
                )
                return
            if operation.get("op") not in ALLOWED_EDIT_OPERATIONS:
                self.send_json(
                    HTTPStatus.BAD_REQUEST,
                    {
                        "error": f"Unsupported operation: {operation.get('op')}",
                        "allowed": sorted(ALLOWED_EDIT_OPERATIONS),
                    },
                )
                return

        job_id = job_dir.name
        try:
            outputs = apply_edit_operations(job_dir, operations)
        except (EditError, IllegalCharacterError, ValueError) as error:
            self.send_json(HTTPStatus.BAD_REQUEST, {"error": str(error)})
            return
        except OSError as error:
            self.send_json(
                HTTPStatus.INTERNAL_SERVER_ERROR,
                {"error": "Edit failed due to a filesystem error"},
            )
            return

        self.send_json(
            HTTPStatus.OK,
            {
                "jobId": job_id,
                "applied": len(operations),
                "files": [file_metadata(path, job_id) for path in outputs],
            },
        )

    def handle_generate(self) -> None:
        payload = self._read_json_body()
        if payload is None:
            return

        title = payload.get("title")
        content = payload.get("content")
        if not isinstance(title, str) or not title.strip():
            self.send_json(HTTPStatus.BAD_REQUEST, {"error": "title is required"})
            return
        if not isinstance(content, str) or not content.strip():
            self.send_json(HTTPStatus.BAD_REQUEST, {"error": "content is required"})
            return
        if len(title) > MAX_TITLE_CHARS:
            self.send_json(
                HTTPStatus.BAD_REQUEST,
                {"error": f"title must be at most {MAX_TITLE_CHARS} characters"},
            )
            return
        if len(content) > MAX_CONTENT_CHARS:
            self.send_json(
                HTTPStatus.BAD_REQUEST,
                {"error": f"content must be at most {MAX_CONTENT_CHARS} characters"},
            )
            return
        if ILLEGAL_XML_CONTROL.search(title) or ILLEGAL_XML_CONTROL.search(content):
            self.send_json(
                HTTPStatus.BAD_REQUEST,
                {"error": "title or content contains unsupported control characters"},
            )
            return

        job_id = uuid.uuid4().hex
        job_dir = BASE_DIR / job_id
        try:
            with JOB_LOCK:
                cleanup_expired_jobs()
                job_directories = [
                    path for path in BASE_DIR.iterdir() if path.is_dir()
                ]
                if len(job_directories) >= MAX_JOBS:
                    self.send_json(
                        HTTPStatus.TOO_MANY_REQUESTS,
                        {"error": "Session job limit reached"},
                    )
                    return
                if directory_size(BASE_DIR) >= MAX_STORAGE_BYTES:
                    self.send_json(
                        HTTPStatus.INSUFFICIENT_STORAGE,
                        {"error": "Session storage limit reached"},
                    )
                    return
                job_dir.mkdir(parents=True)
        except OSError as error:
            self.send_json(
                HTTPStatus.INTERNAL_SERVER_ERROR,
                {"error": "Unable to allocate job workspace"},
            )
            return

        docx_path = job_dir / "report.docx"
        pdf_path = job_dir / "report.pdf"
        pptx_path = job_dir / "report.pptx"
        xlsx_path = job_dir / "report.xlsx"

        try:
            document = Document()
            document.add_heading(title.strip(), level=0)
            for line in content.strip().splitlines():
                document.add_paragraph(line)
            document.save(docx_path)

            workbook = Workbook()
            worksheet = workbook.active
            worksheet.title = "Report"
            worksheet["A1"] = spreadsheet_literal(title.strip())
            worksheet["A2"] = "Generated by AI Workspace Office Sandbox"
            row_index = 4
            for line in content.strip().splitlines():
                chunks = [
                    line[start : start + EXCEL_CELL_MAX_CHARS]
                    for start in range(0, len(line), EXCEL_CELL_MAX_CHARS)
                ] or [""]
                for chunk in chunks:
                    worksheet.cell(
                        row=row_index,
                        column=1,
                        value=spreadsheet_literal(chunk),
                    )
                    row_index += 1
            worksheet.column_dimensions["A"].width = 80
            workbook.save(xlsx_path)

            presentation = Presentation()
            title_slide = presentation.slides.add_slide(
                presentation.slide_layouts[0]
            )
            title_slide.shapes.title.text = title.strip()
            title_slide.placeholders[1].text = "Generated by AI Workspace Office Sandbox"
            content_slide = presentation.slides.add_slide(
                presentation.slide_layouts[1]
            )
            content_slide.shapes.title.text = "Summary"
            text_frame = content_slide.placeholders[1].text_frame
            text_frame.clear()
            for index, line in enumerate(content.strip().splitlines()):
                paragraph = (
                    text_frame.paragraphs[0]
                    if index == 0
                    else text_frame.add_paragraph()
                )
                paragraph.text = line
            presentation.save(pptx_path)

            subprocess.run(
                [
                    "libreoffice",
                    f"-env:UserInstallation=file:///tmp/lo-profile-{job_id}",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(job_dir),
                    str(docx_path),
                ],
                check=True,
                capture_output=True,
                text=True,
                timeout=120,
            )
        except subprocess.TimeoutExpired as error:
            remove_job(job_dir)
            self.send_json(
                HTTPStatus.GATEWAY_TIMEOUT,
                {"error": f"Document conversion timed out: {error.cmd[0]}"},
            )
            return
        except subprocess.CalledProcessError as error:
            remove_job(job_dir)
            self.send_json(
                HTTPStatus.INTERNAL_SERVER_ERROR,
                {
                    "error": f"Document conversion failed: {error.cmd[0]}",
                    "stderr": (error.stderr or "")[-2000:],
                },
            )
            return
        except (IllegalCharacterError, KeyError, ValueError):
            remove_job(job_dir)
            self.send_json(
                HTTPStatus.INTERNAL_SERVER_ERROR,
                {"error": "Document generation failed"},
            )
            return
        except OSError:
            remove_job(job_dir)
            self.send_json(
                HTTPStatus.INTERNAL_SERVER_ERROR,
                {"error": "Document tool could not start"},
            )
            return

        if not all(
            path.is_file()
            for path in (docx_path, pdf_path, pptx_path, xlsx_path)
        ):
            remove_job(job_dir)
            self.send_json(
                HTTPStatus.INTERNAL_SERVER_ERROR,
                {"error": "Expected document artifacts were not generated"},
            )
            return

        try:
            with JOB_LOCK:
                if directory_size(BASE_DIR) > MAX_STORAGE_BYTES:
                    remove_job(job_dir)
                    self.send_json(
                        HTTPStatus.INSUFFICIENT_STORAGE,
                        {"error": "Generated artifacts exceed the session storage limit"},
                    )
                    return
        except OSError:
            remove_job(job_dir)
            self.send_json(
                HTTPStatus.INTERNAL_SERVER_ERROR,
                {"error": "Unable to validate generated artifacts"},
            )
            return

        try:
            files = [
                file_metadata(docx_path, job_id),
                file_metadata(pdf_path, job_id),
                file_metadata(pptx_path, job_id),
                file_metadata(xlsx_path, job_id),
            ]
        except OSError:
            remove_job(job_dir)
            self.send_json(
                HTTPStatus.INTERNAL_SERVER_ERROR,
                {"error": "Unable to read generated artifact metadata"},
            )
            return

        self.send_json(HTTPStatus.OK, {"jobId": job_id, "files": files})

    def log_message(self, format_string: str, *args: object) -> None:
        status = args[1] if len(args) > 1 else None
        print(
            json.dumps(
                {
                    "client": self.client_address[0],
                    "method": self.command,
                    "path": urlparse(self.path).path,
                    "status": status,
                }
            ),
            flush=True,
        )


if __name__ == "__main__":
    BASE_DIR.mkdir(parents=True, exist_ok=True)
    ThreadingHTTPServer(("0.0.0.0", PORT), Handler).serve_forever()
